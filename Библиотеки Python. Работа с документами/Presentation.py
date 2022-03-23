import pptx

sample = '''sample(population, k, *, counts=None) method of random.Random instance
    Chooses k unique random elements from a population sequence or set.

    Returns a new list containing elements from the population while
    leaving the original population unchanged.  The resulting list is
    in selection order so that all sub-slices will also be valid random
    samples.  This allows raffle winners (the sample) to be partitioned
    into grand prize and second place winners (the subslices).

    Members of the population need not be hashable or unique.  If the
    population contains repeats, then each occurrence is a possible
    selection in the sample.

    Repeated elements can be specified one at a time or with the optional
    counts parameter.  For example:

        sample(['red', 'blue'], counts=[4, 2], k=5)

    is equivalent to:

        sample(['red', 'red', 'red', 'red', 'blue', 'blue'], k=5)

    To choose a sample from a range of integers, use range() for the
    population argument.  This is especially fast and space efficient
    for sampling from a large population:

        sample(range(10000000), 60)'''

choice = '''choice(seq) method of random.Random instance
    Choose a random element from a non-empty sequence.'''

random = '''random() -> x in the interval [0, 1).'''

shuffle = '''shuffle(x, random=None) method of random.Random instance
    Shuffle list x in place, and return None.

    Optional argument random is a 0-argument function returning a
    random float in [0.0, 1.0); if it is the default None, the
    standard random.random will be used.'''

gauss = '''gauss(mu, sigma) method of random.Random instance
    Gaussian distribution.

    mu is the mean, and sigma is the standard deviation.  This is
    slightly faster than the normalvariate() function.

    Not thread-safe without a lock around calls.'''

# You can use your own paragraphs and add them in dictionary named 'd'.

prs = pptx.Presentation()
d = {
    sample: 'sample',
    choice: 'choice',
    random: 'random',
    shuffle: 'shuffle',
    gauss: 'gauss'
}
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'Некоторые функции модуля random'
subtitle.text = 'Презентация выполнена Молчановым Степаном'
for i in (sample, choice, random, shuffle, gauss):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = d[i]
    paragraph = slide.shapes.placeholders[1].text_frame.paragraphs[0]
    paragraph.text = i
    paragraph.font.name = 'Courier New'
    paragraph.font.color.rgb = pptx.dml.color.RGBColor(47, 53, 59)
    paragraph.font.size = pptx.util.Pt(12)
    paragraph.alignment = pptx.enum.text.PP_ALIGN.LEFT
prs.save('res.pptx')
